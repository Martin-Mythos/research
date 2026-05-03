from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

slides = [
    ("2026 AI Agent 企业级落地实践与安全挑战", "从 PoC 到规模化：价值闭环、治理框架与风险控制"),
    ("1. 企业落地现状：从工具试点到流程重构", "• 高频场景：客服、知识检索、销售助理、研发 Copilot\n• 瓶颈转移：数据与流程割裂\n• 问题：ROI 难量化、责任边界不清晰"),
    ("2. 参考架构：AgentOps 四层模型", "1) 交互层\n2) 编排层\n3) 治理层\n4) 基础层"),
    ("3. 安全挑战：四类高风险面", "提示词注入 / 数据污染 / 行动失控 / 合规缺失\n对策：隔离、签名、闸门、审计"),
    ("4. 落地路径：90 天推进", "0-30 天：盘点与基线\n31-60 天：双轨试点\n61-90 天：复盘与复制"),
    ("5. KPI 与董事会沟通口径", "效率、质量、风险三类指标并行跟踪")
]

for title, body in slides:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.text = body
    for p in tf.paragraphs:
        p.font.size = Pt(22)

prs.save("artifact-huashu-output.pptx")
print("已生成 artifact-huashu-output.pptx")
